"""Phase 6 工具测试(需求 3):extract_pdf/docx/pptx + transcribe_audio + 项目写操作 + ask_user。"""

import json
import subprocess
import threading

import pytest

from epictrace.config import AppConfig
from epictrace.db import Database
from epictrace.cowork.approvals import ApprovalManager
from epictrace.cowork.tools.builtin_ask import build_ask_user_tool
from epictrace.cowork.tools.builtin_extract import build_extract_tools, build_transcribe_tool
from epictrace.cowork.tools.builtin_projects import build_project_tools
from epictrace.cowork.tools.registry import ToolRegistry


@pytest.fixture()
def env(tmp_path):
    """一个真实项目文件夹 + 注册了 extract/transcribe/project 工具的 registry。"""
    config = AppConfig(data_dir=tmp_path)
    db = Database(config)
    db.create_all()
    proj_dir = tmp_path / "proj"
    proj_dir.mkdir()
    from epictrace.services.projects import ProjectService

    proj = ProjectService(db).create(title="测试项目", folder_path=str(proj_dir))
    registry = ToolRegistry()
    for t in build_extract_tools(db):
        registry.register(t)
    registry.register(build_transcribe_tool(db, is_asr_ready=lambda: True))
    for t in build_project_tools(
        db, get_embedder=lambda: None, get_vector_store=lambda: None,
        index_jobs={}, index_lock=threading.Lock(),
    ):
        registry.register(t)
    return {"db": db, "proj": proj, "dir": proj_dir, "registry": registry}


# ---- extract_pdf/docx/pptx ----

def test_extract_docx(env):
    import docx as docxlib

    doc = docxlib.Document()
    doc.add_paragraph("合同关键条款:违约金按日万分之五")
    doc.save(str(env["dir"] / "contract.docx"))
    out = env["registry"].execute("extract_docx", json.dumps(
        {"project_id": env["proj"].id, "path": "contract.docx"}))
    assert "违约金" in out


def test_extract_pptx(env):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "季度营收汇报"
    prs.save(str(env["dir"] / "q3.pptx"))
    out = env["registry"].execute("extract_pptx", json.dumps(
        {"project_id": env["proj"].id, "path": "q3.pptx"}))
    assert "季度营收汇报" in out


def test_extract_wrong_extension_and_missing(env):
    (env["dir"] / "note.txt").write_text("纯文本", encoding="utf-8")
    out = env["registry"].execute("extract_pdf", json.dumps(
        {"project_id": env["proj"].id, "path": "note.txt"}))
    assert out.startswith("Error:")
    out = env["registry"].execute("extract_docx", json.dumps(
        {"project_id": env["proj"].id, "path": "ghost.docx"}))
    assert "not found" in out
    out = env["registry"].execute("extract_pdf", json.dumps(
        {"project_id": 9999, "path": "x.pdf"}))
    assert "not found" in out


def test_extract_path_escape_blocked(env):
    out = env["registry"].execute("extract_pdf", json.dumps(
        {"project_id": env["proj"].id, "path": "../../etc/passwd"}))
    assert out.startswith("Error:")


# ---- transcribe_audio ----

def test_transcribe_not_ready(env):
    registry = ToolRegistry()
    registry.register(build_transcribe_tool(env["db"], is_asr_ready=lambda: False))
    (env["dir"] / "a.wav").write_bytes(b"RIFF")
    out = registry.execute("transcribe_audio", json.dumps(
        {"project_id": env["proj"].id, "path": "a.wav"}))
    assert "ASR 模型未就绪" in out


def test_transcribe_formats_segments(env, monkeypatch):
    """子进程产出 JSON → 带时间戳文本(不真跑模型,monkeypatch subprocess.run)。"""
    import epictrace.cowork.tools.builtin_extract as mod

    fake = subprocess.CompletedProcess(
        args=[], returncode=0,
        stdout=json.dumps({"segments": [
            {"text": "大家好", "start": 0.0, "end": 1.2, "words": []},
            {"text": "今天讨论预算", "start": 65.0, "end": 68.0, "words": []},
        ]}), stderr="")
    monkeypatch.setattr(mod.subprocess, "run", lambda *a, **k: fake)
    (env["dir"] / "meeting.wav").write_bytes(b"RIFF")
    out = env["registry"].execute("transcribe_audio", json.dumps(
        {"project_id": env["proj"].id, "path": "meeting.wav"}))
    assert "[00:00] 大家好" in out
    assert "[01:05] 今天讨论预算" in out


def test_transcribe_bad_extension(env):
    (env["dir"] / "x.txt").write_text("t", encoding="utf-8")
    out = env["registry"].execute("transcribe_audio", json.dumps(
        {"project_id": env["proj"].id, "path": "x.txt"}))
    assert "不是支持的音频格式" in out


# ---- 项目写操作 ----

def test_create_project(env):
    out = env["registry"].execute("create_project", '{"title": "新项目 Alpha"}')
    assert "新项目 Alpha" in out and "id=" in out
    from epictrace.services.projects import ProjectService

    titles = [p.title for p in ProjectService(env["db"]).list()]
    assert "新项目 Alpha" in titles


def test_add_file_to_project(env):
    src = env["dir"].parent / "outside.txt"
    src.write_text("外部文件内容", encoding="utf-8")
    out = env["registry"].execute("add_file_to_project", json.dumps(
        {"project_id": env["proj"].id, "source_path": str(src)}))
    assert "已添加" in out
    # 文件复制进了项目文件夹
    assert (env["dir"] / "outside.txt").exists()


def test_add_file_missing(env):
    out = env["registry"].execute("add_file_to_project", json.dumps(
        {"project_id": env["proj"].id, "source_path": "/nonexistent/x.pdf"}))
    assert "不存在" in out


def test_rebuild_index_reuses_lock_and_jobs(env):
    from tests.fakes import FakeEmbedder, FakeVectorStore

    jobs = {}
    lock = threading.Lock()
    registry = ToolRegistry()
    for t in build_project_tools(
        env["db"], get_embedder=lambda: FakeEmbedder(),
        get_vector_store=lambda: FakeVectorStore(),
        index_jobs=jobs, index_lock=lock,
    ):
        registry.register(t)
    out = registry.execute("rebuild_index", json.dumps({"project_id": env["proj"].id}))
    assert "启动索引" in out
    assert env["proj"].id in jobs
    job = jobs[env["proj"].id]
    # 后台线程跑完(空项目 total=0 即刻完成):状态机到达 done
    for _ in range(100):
        if job.status != "running":
            break
        import time
        time.sleep(0.05)
    assert job.status == "done"


def test_project_tools_require_confirmation(env):
    for name in ("create_project", "add_file_to_project", "rebuild_index"):
        assert env["registry"].get(name).permission == "ask"


# ---- ask_user ----

def test_ask_user_roundtrip():
    mgr = ApprovalManager()
    tool = build_ask_user_tool(mgr, session_id=7)
    assert tool.permission == "allow"

    holder = {}

    def call():
        holder["out"] = tool.handler("你偏好哪种格式?")

    t = threading.Thread(target=call)
    t.start()
    import time
    for _ in range(50):
        pending = mgr.pending()
        if pending:
            break
        time.sleep(0.05)
    assert pending[0]["kind"] == "question"
    assert pending[0]["prompt"] == "你偏好哪种格式?"
    mgr.decide(pending[0]["approval_id"], "Markdown 就好")
    t.join(timeout=5)
    assert holder["out"] == "用户的回答:Markdown 就好"


def test_ask_user_timeout_returns_placeholder(monkeypatch):
    import epictrace.cowork.tools.builtin_ask as mod

    monkeypatch.setattr(mod, "APPROVAL_TIMEOUT_SEC", 0.05)
    tool = build_ask_user_tool(ApprovalManager(), session_id=1)
    assert "未及时回答" in tool.handler("在吗?")


def test_question_decision_validation():
    mgr = ApprovalManager()
    req = mgr.request(session_id=1, tool="ask_user", args="",
                      allow_session_option=False, kind="question", prompt="q")
    with pytest.raises(ValueError):
        mgr.decide(req["approval_id"], "   ")  # 空回答非法
    assert mgr.decide(req["approval_id"], "答") is True
