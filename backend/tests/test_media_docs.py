from pathlib import Path

import pytest

from epictrace.config import AppConfig
from epictrace.media import get_processor
from epictrace.media.docx import DocxMediaProcessor
from epictrace.media.errors import ExtractionEngineNotReady
from epictrace.media.mineru import MinerUMediaProcessor
from epictrace.media.pdf import PdfMediaProcessor
from epictrace.media.pptx import PptxMediaProcessor


def _set_engine(tmp_path: Path, engine: str) -> AppConfig:
    """持久化 extraction engine,返回对应 AppConfig。"""
    from epictrace.services.settings import SettingsService
    cfg = AppConfig(data_dir=tmp_path)
    SettingsService(cfg).set_extraction_settings(
        engine=engine, effort="medium", model_source="modelscope")
    return cfg


# ---- v2: engine=pypdf(默认)→ 富文档走内置 pypdf/python-docx/python-pptx 处理器 ----


def test_default_engine_is_pypdf_for_rich_docs(tmp_path: Path):
    # 无持久化设置 → 默认 pypdf 引擎;pdf/docx/pptx 走内置处理器(免安装、免模型)。
    cfg = AppConfig(data_dir=tmp_path)
    assert isinstance(get_processor(tmp_path / "a.pdf", cfg), PdfMediaProcessor)
    assert isinstance(get_processor(tmp_path / "a.docx", cfg), DocxMediaProcessor)
    assert isinstance(get_processor(tmp_path / "a.pptx", cfg), PptxMediaProcessor)


def test_pypdf_engine_explicit(tmp_path: Path):
    cfg = _set_engine(tmp_path, "pypdf")
    assert isinstance(get_processor(tmp_path / "a.pdf", cfg), PdfMediaProcessor)
    assert isinstance(get_processor(tmp_path / "a.docx", cfg), DocxMediaProcessor)
    assert isinstance(get_processor(tmp_path / "a.pptx", cfg), PptxMediaProcessor)


def test_pypdf_processors_extract_text(tmp_path: Path):
    # 内置处理器对真实小文件能直接抽文字(无需引擎/模型)。
    import docx as docxlib
    from pptx import Presentation

    cfg = _set_engine(tmp_path, "pypdf")
    d = tmp_path / "hi.docx"
    doc = docxlib.Document()
    doc.add_paragraph("hello docx")
    doc.save(str(d))
    proc = get_processor(d, cfg)
    assert "hello docx" in proc.process(d).text

    pp = tmp_path / "hi.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "hello pptx"
    prs.save(str(pp))
    pproc = get_processor(pp, cfg)
    assert "hello pptx" in pproc.process(pp).text


# ---- v2: engine=mineru → 富文档走 MinerU(无回退) ----


@pytest.mark.parametrize("name", ["a.pdf", "a.docx", "a.pptx"])
def test_mineru_engine_routes_rich_docs_to_mineru(tmp_path: Path, name: str):
    cfg = _set_engine(tmp_path, "mineru")
    p = tmp_path / name
    p.write_bytes(b"x")  # 只验证选路 + 无回退;不需真实文件内容
    proc = get_processor(p, cfg)
    assert isinstance(proc, MinerUMediaProcessor)
    # 未 provision → 处理报错(无回退,不返回 python 处理器文本)
    with pytest.raises(ExtractionEngineNotReady):
        proc.process(p)


def test_mineru_engine_uses_persisted_extraction_settings(tmp_path: Path):
    # 持久化了 mineru + effort=high / model_source=huggingface → registry 据此构造处理器。
    from epictrace.services.settings import SettingsService
    cfg = AppConfig(data_dir=tmp_path)
    SettingsService(cfg).set_extraction_settings(
        engine="mineru", effort="high", model_source="huggingface")
    proc = get_processor(tmp_path / "a.pdf", cfg)
    assert isinstance(proc, MinerUMediaProcessor)
    assert proc._effort == "high"
    assert proc._model_source == "huggingface"


# ---- text/code/data 始终走 TextMediaProcessor,与 engine 无关 ----


@pytest.mark.parametrize("engine", ["pypdf", "mineru"])
def test_text_files_always_text_processor(tmp_path: Path, engine: str):
    from epictrace.media.text import TextMediaProcessor
    cfg = _set_engine(tmp_path, engine)
    assert isinstance(get_processor(tmp_path / "a.md", cfg), TextMediaProcessor)
    assert isinstance(get_processor(tmp_path / "a.py", cfg), TextMediaProcessor)
    assert isinstance(get_processor(tmp_path / "a.csv", cfg), TextMediaProcessor)


@pytest.mark.parametrize("engine", ["pypdf", "mineru"])
def test_unknown_type_returns_none(tmp_path: Path, engine: str):
    cfg = _set_engine(tmp_path, engine)
    assert get_processor(tmp_path / "a.png", cfg) is None    # 图片本期无 processor
    assert get_processor(tmp_path / "a.mp3", cfg) is None     # 音频本期无 processor
