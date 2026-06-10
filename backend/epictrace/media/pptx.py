from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from epictrace.interfaces.media import MediaProcessor, MediaResult


class PptxMediaProcessor(MediaProcessor):
    def supports(self, path: Path) -> bool:
        return path.suffix.lower() == ".pptx"

    def process(self, path: Path) -> MediaResult:
        prs = Presentation(str(path))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        run_text = "".join(run.text for run in para.runs) or para.text
                        if run_text:
                            lines.append(run_text)
        return MediaResult(text="\n".join(lines), metadata={"slides": len(prs.slides)})
